#!/usr/bin/env python3
"""
Create PowerPoint presentation for STI: Making AI Work for All
Professional theme with compelling visuals
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# Color scheme - Professional teal/navy with warm accents
COLORS = {
    'primary': RGBColor(0, 95, 115),      # Deep teal
    'secondary': RGBColor(10, 36, 99),     # Navy blue
    'accent': RGBColor(255, 166, 43),      # Warm amber
    'light': RGBColor(240, 248, 255),      # Light background
    'text_dark': RGBColor(33, 37, 41),     # Near black
    'text_light': RGBColor(255, 255, 255), # White
    'highlight': RGBColor(220, 53, 69),    # Red for emphasis
}

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs)

    # Slide 2: Core Question
    add_content_slide(prs,
        "The Core Question",
        [
            "Will we manage AI's challenges better than we managed globalization?",
            "",
            "• We are facing a collapse in the cost of intelligence",
            "• Technology spreads instantly; prosperity does not",
            "• The Thesis: The central challenge is not the technology itself, but the adoption architecture required to convert disruption into shared prosperity"
        ],
        quote="Will we repeat the mistakes of the past, or write a new chapter?",
        image_suggestion="Split image: Globalization protests on left, AI future on right")

    # Slide 3: Agenda
    add_content_slide(prs,
        "Agenda",
        [
            "1. The Landscape: Why this time is different",
            "   Speed, Cost, Agency",
            "",
            "2. The Evidence: Live demonstrations of AI in policy workflows",
            "",
            "3. The History: Lessons from Industrial Revolution & IT Era",
            "",
            "4. The Challenges: Stability, Labor, Explainability, Adoption",
            "",
            "5. The Strategy: A framework for national adoption"
        ],
        image_suggestion="Abstract timeline or roadmap graphic")

    # Slide 4: Opening Provocation - Ricardo
    add_content_slide(prs,
        "Opening Provocation: The Ricardo Pivot (1821)",
        [
            "• 1817: David Ricardo argues machinery benefits everyone",
            "",
            "• 1821: After seeing wages collapse, he changes his mind:",
            '  "Machinery may be detrimental to the interests of the class of laborers"',
            "",
            "• Today: We have the benefit of hindsight",
            "  Technology creates wealth, but policy determines distribution"
        ],
        quote="Technology is inevitable. Inclusive growth is a choice.",
        image_suggestion="Historical illustration of early industrial machinery or Ricardo portrait")

    # Slide 5: Speed of Diffusion
    add_content_slide(prs,
        "Context: The Speed of Diffusion",
        [
            "Time to 100 Million Users:",
            "",
            "• Internet: 7 Years",
            "• Facebook: 4.5 Years",
            "• ChatGPT: 2 Months",
            "",
            "Implication: The 'Policy Response Window' has collapsed",
            "Traditional multi-year regulatory cycles cannot keep pace with monthly capability jumps"
        ],
        image_suggestion="Bar chart showing adoption timelines - visually striking comparison")

    # Slide 6: New Economics
    add_content_slide(prs,
        "Context: The Collapse of Cognitive Costs",
        [
            "Zero Marginal Cost:",
            "• Cost of summarizing 100-page report: $100s (analyst) → $0.01 (compute)",
            "",
            "Cognitive Automation:",
            "• For the first time, automation targets knowledge work",
            "  Analysis, writing, coding - not just manual labor",
            "",
            "The Shift:",
            "• From 'Tools' (Excel) to 'Agents' (Systems that plan and execute)"
        ],
        image_suggestion="Cost curve dropping exponentially or brain/circuit hybrid image")

    # Slide 7: Horizon Scan
    add_content_slide(prs,
        "The Horizon Scan: Where Are We Going?",
        [
            "Horizon 1 (Now): Copilots",
            "• Human-in-the-loop | Augmentation | 20-40% productivity gain",
            "",
            "Horizon 2 (2026+): Agents",
            "• Human-on-the-loop | Delegation | Workflow transformation",
            "",
            "Horizon 3 (2028+): System Reshuffle",
            "• 'Coordination without Consensus' | Dynamic automated teams"
        ],
        image_suggestion="Three-stage timeline with icons: person+AI, person supervising AI, interconnected AI systems")

    # Slide 8: Demo Introduction
    add_section_slide(prs, "Live Demonstrations", "Seeing is Believing",
        subtitle="AI is a reasoning engine, not a knowledge database")

    # Slide 9: Demo 1
    add_content_slide(prs,
        "Demo 1: Policy Drafting",
        [
            "Task: Draft a briefing note on 'Healthcare Policy for Aging Population'",
            "",
            "Prompt: 'Act as a senior policy analyst. Draft a briefing note on integrating AI into eldercare. Cover: Privacy risks, workforce augmentation, infrastructure needs.'",
            "",
            "Result: A structured, nuanced draft in 30 seconds",
            "",
            "Implication: The 'Blank Page Problem' is solved",
            "Policymakers shift from writers to editors and strategic thinkers"
        ],
        image_suggestion="Screenshot of AI generating policy document or before/after comparison")

    # Slide 10: Demo 2
    add_content_slide(prs,
        "Demo 2: Synthesis at Scale",
        [
            "Task: Analyze Public Consultation Feedback",
            "",
            "Scenario: 500 emails from citizens on a new transport subsidy",
            "",
            "Prompt: 'Analyze these consultation responses. Identify top 3 concerns, prevailing sentiment, suggest 2 mitigation strategies.'",
            "",
            "Result: Instant thematic analysis of unstructured text",
            "",
            "Implication: Democracy becomes more responsive",
            "Process qualitative feedback at the scale of quantitative data"
        ],
        image_suggestion="Word cloud or sentiment analysis visualization")

    # Slide 11: Demo 3
    add_content_slide(prs,
        "Demo 3: Red Teaming Policy",
        [
            "Task: Stress-test a proposed regulation",
            "",
            "Scenario: New policy to tax digital services",
            "",
            "Prompt: 'You are the Opposition Leader. Red team this policy. Find loopholes, unintended consequences, implementation blind spots.'",
            "",
            "Result: Rigorous, adversarial critique instantly",
            "",
            "Implication: Stronger policy-making through 'AI Adversaries'"
        ],
        image_suggestion="Chess pieces or debate/adversarial imagery")

    # Slide 12: Historical Lesson 1
    add_content_slide(prs,
        "Historical Lesson 1: Engels' Pause",
        [
            "The Data (Industrial Revolution):",
            "• 1780-1840: Productivity rose 46%",
            "• But wages rose only 12%",
            "",
            "The Gap: 60 years for workers to feel the benefits",
            "",
            "Why? Institutions lagged behind technology",
            "• Unions, Education, Safety Nets",
            "",
            "Our Goal: Shorten the pause"
        ],
        image_suggestion="Graph showing productivity vs wages divergence, or historical factory image")

    # Slide 13: Historical Lesson 2
    add_content_slide(prs,
        "Historical Lesson 2: The Solow Paradox",
        [
            '"You can see the computer age everywhere but in the productivity statistics."',
            "— Robert Solow, 1987",
            "",
            "The Lesson:",
            "Electrifying a steam factory didn't help until the workflow was redesigned",
            "",
            "Application:",
            "Paving cow paths (digitizing bad processes) yields zero gain",
            "We must reimagine public services, not just add AI to them"
        ],
        image_suggestion="Factory floor evolution or computer in traditional office")

    # Slide 14: Perez Framework Intro
    add_section_slide(prs, "Understanding Technological Revolutions",
        "The Perez Framework",
        subtitle="Five Great Surges Since 1771")

    # Slide 15: Five Surges
    add_content_slide(prs,
        "The Perez Framework: Five Great Surges",
        [
            "1. Industrial Revolution (1771)",
            "2. Age of Steam & Railways (1829)",
            "3. Age of Steel & Electricity (1875)",
            "4. Age of Oil, Automobiles & Mass Production (1908)",
            "5. Age of Information & Telecommunications (1971)",
            "",
            "The Claim: Each follows a remarkably similar pattern",
            "Diffusion → Crisis → 'Golden Age'",
            "",
            "AI (2022?): We may be witnessing the sixth"
        ],
        image_suggestion="Timeline with icons for each revolution - visual history")

    # Slide 16: Four Phases
    add_content_slide(prs,
        "The Four Phases of a Revolution",
        [
            "1. Irruption (0-10 years)",
            "   New technology emerges. Financial capital floods in.",
            "",
            "2. Frenzy (10-20 years) ← WE ARE LIKELY HERE",
            "   Speculation explodes. Infrastructure overbuilt.",
            "",
            "3. Turning Point",
            "   Bubble bursts. Regulatory reckoning begins.",
            "",
            "4. Synergy/Deployment (20-40 years)",
            "   Technology spreads broadly. 'Golden Ages' emerge."
        ],
        image_suggestion="S-curve or wave diagram with phases marked")

    # Slide 17: Bubbles as Features
    add_content_slide(prs,
        "Why Bubbles Are a Feature, Not a Bug",
        [
            "The Paradox: Bubbles are destructive, yet essential",
            "",
            "The Function: Irrational exuberance mobilizes capital at scale",
            "",
            "• Railways wildly overbuilt in 1840s → but we got infrastructure",
            "• Fiber optics over-deployed in 1990s → but we got cheap bandwidth",
            "• AI data centers being built at unprecedented scale → ...",
            "",
            "The Lesson: Not 'will there be a crash?'",
            "But 'what infrastructure will remain when the music stops?'"
        ],
        image_suggestion="Bubble imagery or infrastructure being built")

    # Slide 18: Turning Point
    add_content_slide(prs,
        "The Turning Point: What Happens Next?",
        [
            "The Crisis: Financial crashes (1847, 1929, 2000) mark the end",
            "",
            "The Reckoning: Society demands regulation",
            "Excesses of 'Gilded Ages' give way to reform",
            "",
            "The Opportunity: This is when policy matters most",
            "Institutional choices determine next 30 years",
            "",
            "For AI: We have a narrow window to shape the framework",
            "Before deployment locks in"
        ],
        image_suggestion="Crossroads or pivot point imagery")

    # Slide 19: Financial Stability
    add_content_slide(prs,
        "Challenge 1: Financial Stability",
        [
            "Systemic Risk & Procyclicality",
            "",
            "Herding: Same models → Same risks identified → Same reactions",
            "",
            "The Result: Massive, synchronized market movements",
            "",
            "The Perez Warning: In Frenzy phase, financial capital dominates",
            "It seeks returns, not deployment. Prepare for correction.",
            "",
            "Policy Action: Demand 'Model Diversity'",
            "Stress-test for algorithmic correlation"
        ],
        image_suggestion="Synchronized dominoes or herding visualization")

    # Slide 20: Labor Markets Overview
    add_section_slide(prs, "Challenge 2: Labor Markets",
        "Jobs, Tasks, and the Dual Effect",
        subtitle="Displacement vs. Reinstatement")

    # Slide 21: Dual Effect
    add_content_slide(prs,
        "The Dual Effect of Technology",
        [
            "Displacement Effect:",
            "• Automation replaces labor in existing tasks",
            "",
            "Reinstatement Effect:",
            "• New technology creates entirely new tasks and jobs",
            "",
            "Historical Record: Both effects always occur",
            "The question is timing and distribution",
            "",
            "The AI Difference: First time automation targets",
            "cognitive tasks at scale — analysis, writing, coding, diagnosis"
        ],
        image_suggestion="Balance scale or two-way arrows diagram")

    # Slide 22: Task-Based Framework
    add_content_slide(prs,
        "Thinking in Tasks, Not Jobs",
        [
            "The Insight (Acemoglu & Restrepo):",
            "Jobs are bundles of tasks. Technology automates specific tasks.",
            "",
            "Example — Policy Analyst:",
            "• Tasks automated: Literature review, first drafts, data summarization",
            "• Tasks augmented: Strategic synthesis, stakeholder judgment",
            "• Tasks created: Prompt engineering, AI oversight, algorithmic accountability",
            "",
            "The Implication: Job persists, composition changes",
            "Value shifts from execution to judgment"
        ],
        image_suggestion="Job decomposed into task components - visual breakdown")

    # Slide 23: Displacement Concern
    add_content_slide(prs,
        "The Displacement Concern: Who Is Most Exposed?",
        [
            "Routine Cognitive Work: Rule-based, predictable, information-intensive",
            "• Back-office processing, basic legal research, financial analysis",
            "• Customer service, coding",
            "",
            "The Skill Paradox:",
            "AI affects middle and upper-middle of skill distribution",
            "— precisely the jobs that absorbed displaced manufacturing workers",
            "",
            "Geographic & Demographic Concentration:",
            "Displacement will not be evenly distributed"
        ],
        image_suggestion="Skills distribution curve with exposure highlighted")

    # Slide 24: Compensating Effects
    add_content_slide(prs,
        "The Compensating Effects",
        [
            "Productivity Complementarity:",
            "• AI augments rather than replaces — 30-50% productivity gains",
            "",
            "Cost Reduction → Demand Expansion:",
            "• Cheaper cognitive services = more access (legal, education, healthcare)",
            "",
            "New Task Creation:",
            "• Every revolution creates unimaginable tasks",
            "",
            "The Policy Lever: Reinstatement is not automatic",
            "Requires investment in skills, institutions, conditions"
        ],
        image_suggestion="Growth/expansion imagery or new job roles emerging")

    # Slide 25: Experience Trap
    add_content_slide(prs,
        "The Experience Trap: Cutting the Bottom Rungs",
        [
            "The Mechanism:",
            "AI targets routine cognitive tasks — the tasks juniors use to learn",
            "",
            "The Problem:",
            "If AI writes first drafts and does basic analysis,",
            "how do juniors develop tacit knowledge?",
            "",
            "The Risk: High productivity today,",
            "crisis of expertise in 5-10 years",
            "",
            "Examples: Junior lawyers, analysts, doctors, accountants, engineers"
        ],
        image_suggestion="Ladder with bottom rungs removed or pipeline drying up")

    # Slide 26: Skills in AI Era
    add_content_slide(prs,
        "Skills and Expertise in the AI Era",
        [
            "What Becomes Valuable?",
            "",
            "• Tacit Knowledge: Judgment, intuition, pattern recognition",
            "  Cannot be easily codified",
            "",
            "• Visible Expertise: Demonstrating why decisions were made",
            "  Ethics, context, stakeholder understanding",
            "",
            "• Orchestration Skills: Managing AI systems, knowing when to trust",
            "",
            "• Learning Velocity: Ability to unlearn and relearn monthly"
        ],
        image_suggestion="Human-AI collaboration or expertise pyramid")

    # Slide 27: Rethinking Development
    add_content_slide(prs,
        "Rethinking Professional Development",
        [
            "The Old Model:",
            "Learn fundamentals → Practice routine → Build to complexity (years)",
            "",
            "The New Model:",
            "Engage complexity early, AI handles routine execution",
            "Juniors learn to evaluate and direct, not just execute",
            "",
            "Institutional Response Required:",
            "• Restructure apprenticeship and training",
            "• Create 'AI-assisted learning' pathways",
            "• Measure outcomes (judgment, quality) not just throughput"
        ],
        image_suggestion="New learning pathway or mentorship imagery")

    # Slide 28: Explainability 1
    add_content_slide(prs,
        "Challenge 3: The Black Box Dilemma",
        [
            "The Conflict:",
            "• Policy requires accountability and reason",
            "• AI operates on probability and pattern matching, not logic",
            "",
            "The 'Explainability Trilemma':",
            "Can usually only have two: Accuracy, Complexity, Explainability",
            "High-performance models are often least explainable",
            "",
            "The Risk:",
            "'Computer says no' is not acceptable for denied benefits or licenses"
        ],
        image_suggestion="Black box or neural network visualization")

    # Slide 29: Explainability 2
    add_content_slide(prs,
        "Operationalizing Trust",
        [
            "From 'No' to 'Guardrails': Cannot wait for perfect explainability",
            "",
            "Solutions:",
            "",
            "• Sandboxes: Safe spaces with error budgets for experimentation",
            "",
            "• 'AI Verify': Testing frameworks (like Singapore's)",
            "  Benchmark model behavior against safety standards",
            "",
            "• Human-Over-The-Loop:",
            "  Critical decisions must have human accountable for outcome"
        ],
        image_suggestion="Guardrails or safety framework diagram")

    # Slide 30: Adoption Gap 1
    add_content_slide(prs,
        "Challenge 4: The 95% Pilot Trap",
        [
            "The Reality:",
            "Easy to build a cool demo. Incredibly hard to move to production.",
            "",
            "Why?",
            "",
            "• Data Quality: Real-world data is messy, siloed, unstructured",
            "",
            "• Legacy Systems: AI doesn't easily 'talk' to 1990s databases",
            "",
            "• Culture: Fear of replacement leads to resistance"
        ],
        image_suggestion="Pilot to production gap or bridge imagery")

    # Slide 31: Adoption Architecture
    add_content_slide(prs,
        "The Adoption Architecture",
        [
            "The 'Boring' Rails: Success depends on the unsexy stuff",
            "",
            "• Data Governance: Treat data as strategic asset",
            "",
            "• Procurement: Agile contracting, not 5-year waterfall tenders",
            "",
            "• Talent: Upskill broad civil service, not just hire PhDs",
            "",
            "Key Takeaway:",
            "You don't need more 'AI Strategy' documents",
            "You need 'Implementation Playbooks'"
        ],
        image_suggestion="Infrastructure or foundation building imagery")

    # Slide 32: Framework for Action
    add_section_slide(prs, "The Strategy",
        "A Framework for Action",
        subtitle="Three Pillars for National Success")

    # Slide 33: Three Pillars
    add_content_slide(prs,
        "Three Pillars for National Success",
        [
            "1. Institutional Readiness",
            "   Governance, Data, Talent — The Foundation",
            "",
            "2. Strategic Use Cases",
            "   Focus on Utility (Reporting, Services), not Magic",
            "",
            "3. Ecosystem Collaboration",
            "   Research + Government + Private Sector",
            "   The Triple Helix"
        ],
        image_suggestion="Three pillars or triple helix DNA structure")

    # Slide 34: What Governments Must Do
    add_content_slide(prs,
        "What Governments Must Do",
        [
            "1. Fix the Adoption Architecture",
            "   • Modernize procurement for agile AI projects",
            "   • Establish data governance and sharing frameworks",
            "   • Build AI literacy across civil service",
            "",
            "2. Proactive Labor Market Policy",
            "   • Redesign professional development pathways",
            "   • Invest in reinstatement — new tasks, not just protecting old jobs",
            "",
            "3. Whole-of-Government Collaboration",
            "   • Shared infrastructure, common standards, coordinated risk approach"
        ],
        image_suggestion="Government building or policy implementation visual")

    # Slide 35: What Organizations Must Do
    add_content_slide(prs,
        "What Organizations Must Do",
        [
            "1. Redesign Work, Not Just Add AI",
            "   • Audit task composition of key roles",
            "   • Create new roles for AI oversight, ethics, quality",
            "",
            "2. Protect the Experience Pipeline",
            "   • Restructure junior roles to preserve learning",
            "   • Measure development outcomes, not just efficiency",
            "",
            "3. Build Responsible Adoption Culture",
            "   • Train all staff, not just specialists",
            "   • Safe spaces for experimentation"
        ],
        image_suggestion="Organizational transformation or team adaptation")

    # Slide 36: What Can I Do
    add_content_slide(prs,
        "What Can I Do?",
        [
            "Personal Agency in the AI Transition",
            "",
            "• Become AI-Literate: You don't need to code, but understand",
            "  capabilities and limitations. Use the tools. Build intuition.",
            "",
            "• Invest in Complementary Skills: Double down on what AI cannot do",
            "  — judgment, ethics, relationships, context, creativity",
            "",
            "• Develop Learning Velocity: Tools change. Unlearn/relearn is the meta-skill",
            "",
            "• Be a Bridge: Translate between technical possibility and operational reality"
        ],
        image_suggestion="Individual growth or personal development imagery")

    # Slide 37: Practical Steps
    add_content_slide(prs,
        "Practical Steps This Week",
        [
            "1. Experiment: 30 minutes with AI on a real work task",
            "   Draft a brief, analyze feedback, stress-test a proposal",
            "",
            "2. Audit Your Role: Which tasks are routine cognitive?",
            "   Which require judgment, context, relationships?",
            "",
            "3. Have the Conversation: Talk to your team about AI",
            "   Fears? Opportunities? Quick wins?",
            "",
            "4. Champion Pilots: One low-risk, high-visibility use case",
            "",
            "5. Stay Informed: Landscape changes monthly"
        ],
        image_suggestion="Checklist or action steps visual")

    # Slide 38: Individual's Dilemma
    add_content_slide(prs,
        "The Individual's Dilemma",
        [
            "The Easy Path: Wait and see. Let others figure it out.",
            "",
            "The Risk: Those who wait will find the transition",
            "imposed on them, rather than shaped by them",
            "",
            "The Alternative: Engage now, while technology is malleable",
            "Your voice matters more in Installation than Deployment"
        ],
        quote='"In times of change, learners inherit the earth, while the learned find themselves beautifully equipped to deal with a world that no longer exists." — Eric Hoffer',
        image_suggestion="Crossroads or choice imagery")

    # Slide 39: Three Levels Summary
    add_table_slide(prs,
        "Closing: Three Levels of Action",
        [
            ["Level", "Key Action", "Timeline"],
            ["Government", "Fix adoption architecture (data, procurement, talent)", "1-3 years"],
            ["Organization", "Redesign work; protect experience pipeline", "6-18 months"],
            ["Individual", "Build AI literacy; invest in complementary skills", "This week"]
        ],
        footer="Technology is inevitable. Inclusive growth is a choice. That choice starts with you.")

    # Slide 40: Final Thought
    add_final_slide(prs)

    # Save
    output_path = "/home/user/STI_Presentation/STI_Presentation_AI_Work_For_All.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


def add_title_slide(prs):
    """Create compelling title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    add_background(slide, COLORS['secondary'])

    # Accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['accent']
    bar.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Making AI Work for All"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Not Repeating the Mistakes of Globalization"
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER

    # Event info
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1.5))
    tf = info_box.text_frame

    p = tf.paragraphs[0]
    p.text = "A Strategic Briefing for Policymakers"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Singapore Training Institute of the IMF"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "3 December 2025"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Image placeholder note
    add_image_note(slide, "IMAGE: Singapore skyline at dusk or abstract AI visualization", Inches(9.5), Inches(6.5))


def add_section_slide(prs, section_title, main_text, subtitle=None):
    """Create section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    add_background(slide, COLORS['primary'])

    # Section indicator
    ind_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(0.5))
    tf = ind_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title.upper()
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER

    # Main text
    main_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.5))
    tf = main_box.text_frame
    p = tf.paragraphs[0]
    p.text = main_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.color.rgb = COLORS['text_light']
        p.alignment = PP_ALIGN.CENTER


def add_content_slide(prs, title, content_lines, quote=None, image_suggestion=None):
    """Create standard content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Light background
    add_background(slide, COLORS['light'])

    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['secondary']
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']

    # Content
    content_top = Inches(1.5)
    content_height = Inches(4.5) if quote else Inches(5.5)

    content_box = slide.shapes.add_textbox(Inches(0.75), content_top, Inches(11.833), content_height)
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        p.space_after = Pt(6)

        # Style bullets
        if line.strip().startswith('•'):
            p.level = 0
        elif line.strip().startswith('-'):
            p.level = 1

    # Quote if present
    if quote:
        quote_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(1))
        tf = quote_box.text_frame
        p = tf.paragraphs[0]
        p.text = quote
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

    # Image suggestion
    if image_suggestion:
        add_image_note(slide, f"IMAGE: {image_suggestion}", Inches(9.5), Inches(0.3))


def add_table_slide(prs, title, table_data, footer=None):
    """Create slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_background(slide, COLORS['light'])

    # Header
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS['secondary']
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['text_light']

    # Table
    rows = len(table_data)
    cols = len(table_data[0])

    table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(11.333), Inches(2.5)).table

    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text

            # Style header row
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['primary']
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLORS['text_light']
                    paragraph.font.size = Pt(16)
            else:
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.color.rgb = COLORS['text_dark']

    # Footer
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.333), Inches(1))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER


def add_final_slide(prs):
    """Create powerful closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_background(slide, COLORS['secondary'])

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(1.5))
    tf = quote_box.text_frame
    p = tf.paragraphs[0]
    p.text = '"The AI future is unwritten. It is not determined by the algorithm, but by the institutions in this room—and the individuals who lead them."'
    p.font.size = Pt(24)
    p.font.italic = True
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Questions
    q_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.333), Inches(1.2))
    tf = q_box.text_frame

    p = tf.paragraphs[0]
    p.text = "Will we repeat the mistakes of globalization—fast diffusion, slow prosperity?"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Or will we build the adoption architecture for shared benefit?"
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER

    # Call to action
    cta_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.333), Inches(1.5))
    tf = cta_box.text_frame

    p = tf.paragraphs[0]
    p.text = "The choice is made in every policy brief you write, every procurement you approve,"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "every junior you mentor, every conversation you have."
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['text_light']
    p.alignment = PP_ALIGN.CENTER

    # Final question
    final_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.8))
    tf = final_box.text_frame
    p = tf.paragraphs[0]
    p.text = "What will you do Monday morning?"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER


def add_background(slide, color):
    """Add solid color background"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        Inches(13.333),
        Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()

    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_image_note(slide, text, left, top):
    """Add small note for image suggestion"""
    note_box = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.italic = True
    p.font.color.rgb = RGBColor(128, 128, 128)


if __name__ == "__main__":
    create_presentation()
